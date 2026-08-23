import os
import json
from pathlib import Path
import pandas as pd
from fpdf import FPDF
import docx
from pptx import Presentation
from pptx.util import Inches, Pt

# Importa caminhos configurados
try:
    from config import DATA_DIR
except ImportError:
    # Se rodado diretamente
    PROJECT_ROOT = Path(__file__).resolve().parent.parent
    DATA_DIR = PROJECT_ROOT / "data"

DATA_DIR.mkdir(parents=True, exist_ok=True)

def generate_pdf():
    pdf_path = DATA_DIR / "termo_de_uso_e_privacidade.pdf"
    print(f"Gerando PDF: {pdf_path.name}")
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "CloudScale SaaS Solutions - Termo de Uso e Privacidade", ln=True, align="C")
    pdf.ln(10)
    
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, "1. SLA de Servico (Service Level Agreement)", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 6, (
        "A CloudScale SaaS Solutions garante aos seus clientes um SLA de disponibilidade "
        "mensal de 99.9%. Caso esse percentual nao seja atingido, o cliente tera direito "
        "a descontos proporcionais na proxima fatura, conforme previsto no contrato principal."
    ))
    pdf.ln(5)
    
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, "2. Politica de Retencao de Dados e LGPD", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 6, (
        "Em conformidade com a LGPD (Lei Geral de Protecao de Dados), a CloudScale mantem a "
        "politica de retencao de dados de clientes por 5 anos apos o encerramento do contrato. "
        "Durante este periodo, os dados ficam criptografados e acessiveis apenas para fins "
        "de auditoria e cumprimento de obrigacoes legais. Apos 5 anos, os dados sao permanentemente deletados."
    ))
    pdf.ln(5)
    
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, "3. Multas e Penalidades", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 6, (
        "Em caso de descumprimento de obrigacoes contratuais de seguranca por parte do cliente ou "
        "vazamento comprovado por culpa exclusiva, aplicam-se precos de multas de ate 2% do "
        "faturamento anual da empresa infratora, limitado ao valor maximo de R$ 50.000.000 (cinquenta milhoes de reais) "
        "por infracao ocorrente."
    ))
    
    pdf.output(str(pdf_path))

def generate_docx():
    docx_path = DATA_DIR / "manual_onboarding_rh.docx"
    print(f"Gerando DOCX: {docx_path.name}")
    
    doc = docx.Document()
    doc.add_heading("CloudScale SaaS Solutions - Manual de Onboarding de RH", 0)
    
    doc.add_heading("Bem-vindo a Equipe CloudScale!", level=1)
    doc.add_paragraph(
        "Este manual tem como objetivo guiar o novo colaborador no processo de integracao à nossa empresa. "
        "Aqui estao resumidos os principais beneficios e politicas internas."
    )
    
    doc.add_heading("Beneficios Corporativos", level=2)
    doc.add_paragraph("A CloudScale oferece um pacote competitivo de beneficios para todos os colaboradores:")
    
    # Adicionar benefícios em tópicos
    doc.add_paragraph("• Vale Refeicao: Oferecemos Vale Refeicao de R$ 50/dia creditado mensalmente no cartao Flash.", style="List Bullet")
    doc.add_paragraph("• Plano de Saude: Plano de Saude Bradesco Nacional (quarto individual), sem coparticipacao para o titular.", style="List Bullet")
    doc.add_paragraph("• Odontologico: Plano Odontoprev integralmente custeado pela empresa.", style="List Bullet")
    
    doc.add_heading("Politica de Home Office e Trabalho Flexivel", level=2)
    doc.add_paragraph(
        "Adotamos uma Politica de Home Office flexivel baseada em confianca. O colaborador pode optar por trabalhar "
        "no formato 100% remoto ou hibrido (sugerido 2 dias presenciais no escritorio em Sao Paulo para integracao). "
        "A empresa fornece ajuda de custo para internet e equipamentos ergonomicos (cadeira, monitor e notebook de ultima geracao)."
    )
    
    doc.save(str(docx_path))

def generate_xlsx():
    xlsx_path = DATA_DIR / "tabela_planos_e_precos.xlsx"
    print(f"Gerando XLSX: {xlsx_path.name}")
    
    data = {
        "Plano": ["Starter", "Pro", "Enterprise"],
        "Preco Mensal (R$)": [299.00, 799.00, 2499.00],
        "Limite de Requisicoes API/mes": [50000, 250000, "Ilimitado"],
        "Suporte Tecnico": ["Email (24h)", "Chat & Email (12h)", "Gerente Dedicado (24/7)"],
        "Storage Incluido": ["10 GB", "100 GB", "1 TB"]
    }
    
    df = pd.DataFrame(data)
    
    # Criar um writer do pandas usando openpyxl
    with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Planos CloudScale")

def generate_pptx():
    pptx_path = DATA_DIR / "apresentacao_comercial.pptx"
    print(f"Gerando PPTX: {pptx_path.name}")
    
    prs = Presentation()
    
    # Slide 1: Capa
    slide_layout = prs.slide_layouts[0] # layout de capa
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "CloudScale SaaS Solutions"
    subtitle.text = "Apresentacao Comercial - Pitch 2026\nPlataforma de Gestao Multi-Cloud B2B"
    
    # Slide 2: Missao e Metricas
    slide_layout = prs.slide_layouts[1] # layout titulo + conteudo
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Nossa Missao e Metricas de Crescimento"
    
    content_box = shapes.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Missao da Empresa:"
    p = tf.add_paragraph()
    p.text = "Empoderar empresas B2B a gerenciar sua infraestrutura em nuvem de forma escalavel, inteligente e segura, reduzindo custos desnecessarios."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Metricas de Crescimento (2025-2026):"
    p3 = tf.add_paragraph()
    p3.text = "• ARR (Receita Recorrente Anual) atingiu R$ 10M (Dez/2025)"
    p3.level = 1
    p4 = tf.add_paragraph()
    p4.text = "• Mais de 450 clientes corporativos ativos"
    p4.level = 1
    
    # Slide 3: Concorrentes
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Analise de Mercado e Concorrentes"
    content_box = shapes.placeholders[1]
    tf = content_box.text_frame
    tf.text = "Principais Concorrentes Identificados:"
    
    c1 = tf.add_paragraph()
    c1.text = "1. CloudOps Platform: Foco em infraestrutura simples, mas carece de automacao com IA."
    c1.level = 1
    c2 = tf.add_paragraph()
    c2.text = "2. StackManager: Boa integracao com AWS, porem preco elevado para multi-cloud."
    c2.level = 1
    c3 = tf.add_paragraph()
    c3.text = "3. AWS Native Console: Complexo de usar para quem opera com multicloud (Azure e GCP)."
    c3.level = 1
    
    prs.save(str(pptx_path))

def generate_md():
    md_path = DATA_DIR / "faq_suporte_tecnico.md"
    print(f"Gerando MD: {md_path.name}")
    
    content = """# FAQ de Suporte Tecnico - CloudScale SaaS Solutions

Aqui estao as respostas para as principais duvidas tecnicas dos nossos clientes e parceiros de integracao.

## 1. Como resetar tokens de API
Se o seu token de API expirou ou foi exposto, siga o passo a passo para resetar:
1. Acesse o painel administrativo da CloudScale em `console.cloudscale.com`.
2. Navegue ate **Configuracoes da Conta** > **API & Integracoes**.
3. Localize o token atual e clique no botao **Revogar Token**.
4. Em seguida, clique em **Gerar Novo Token**.
5. Copie a nova chave imediatamente, pois ela nao sera exibida novamente.

## 2. Integracao via Webhooks
A CloudScale envia notificacoes em tempo real via Webhooks sobre eventos de billing e alteracoes de status de servidores.
- **Configuracao**: Va para a secao Webhooks no painel e insira a sua URL de recebimento (Endpoint).
- **Seguranca**: Validamos as chamadas adicionando uma assinatura `X-CloudScale-Signature` no header HTTP de cada requisicao.

## 3. Solucao para Erro HTTP 429 (Too Many Requests)
O erro HTTP 429 ocorre quando a sua aplicacao ultrapassa o limite de requisicoes permitidas por minuto para o seu plano:
- **Plano Starter**: Limite maximo de 60 requisicoes por minuto.
- **Plano Pro**: Limite maximo de 300 requisicoes por minuto.
- **Plano Enterprise**: Limite customizado.
**Como mitigar**: Implemente uma estrategia de retentativa com recuo exponencial (Exponential Backoff) e faça cache de dados estaticos localmente para reduzir as chamadas desnecessarias.
"""
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(content)

def generate_csv():
    csv_path = DATA_DIR / "base_clientes_e_retencao.csv"
    print(f"Gerando CSV: {csv_path.name}")
    
    data = {
        "Metric": ["NPS (Net Promoter Score)", "Taxa de Cancelamento (Churn)", "Principais motivos de suporte", "Tempo medio de resposta (SLA)", "Retencao Anual"],
        "Value": ["88", "1.2%", "Duvidas de Integracao API (45%), Billing/Faturamento (30%), Upgrade de Plano (25%)", "12 minutos", "94.5%"],
        "Period": ["Q1 2026", "Q1 2026", "Mensal consolidado", "Q1 2026", "Anual 2025"]
    }
    
    df = pd.DataFrame(data)
    df.to_csv(csv_path, index=False, encoding="utf-8")

def generate_json():
    json_path = DATA_DIR / "configuracao_endpoints_api.json"
    print(f"Gerando JSON: {json_path.name}")
    
    data = {
        "api_name": "CloudScale REST API",
        "version": "v1",
        "base_url": "https://api.cloudscale.com/v1",
        "endpoints": [
            {
                "path": "/v1/auth",
                "method": "POST",
                "description": "Autenticacao do cliente e geracao de JWT temporario.",
                "required_parameters": ["client_id", "client_secret"],
                "optional_parameters": [],
                "error_codes": {
                    "400": "Bad Request - Parametros ausentes",
                    "401": "Unauthorized - Credenciais invalidas"
                }
            },
            {
                "path": "/v1/users",
                "method": "GET",
                "description": "Lista de usuarios administradores da conta cloud.",
                "required_parameters": [],
                "optional_parameters": ["page", "limit"],
                "error_codes": {
                    "401": "Unauthorized - Token expirado",
                    "403": "Forbidden - Permissao insuficiente"
                }
            },
            {
                "path": "/v1/billing",
                "method": "GET",
                "description": "Recupera informacoes de faturamento, faturas abertas e plano ativo.",
                "required_parameters": ["account_id"],
                "optional_parameters": [],
                "error_codes": {
                    "401": "Unauthorized - JWT invalido",
                    "404": "Not Found - Account ID nao encontrado"
                }
            }
        ],
        "rate_limits": {
            "global_default": "100 req/min",
            "error_on_exceed": "HTTP 429 Too Many Requests"
        }
    }
    
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)

def generate_html():
    html_path = DATA_DIR / "politica_seguranca_lgpd.html"
    print(f"Gerando HTML: {html_path.name}")
    
    content = """<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <title>Politica de Seguranca da Informacao - CloudScale</title>
</head>
<body>
    <h1>Politica de Seguranca da Informacao da CloudScale SaaS</h1>
    
    <h2>1. Regras de Criacao de Senhas Fortes</h2>
    <p>Todos os colaboradores e clientes que acessam a plataforma CloudScale devem adotar senhas complexas e seguras. As regras obrigatorias sao:</p>
    <ul>
        <li>Comprimento minimo de <strong>12 caracteres</strong>.</li>
        <li>Conter pelo menos uma letra maiuscula e uma minuscula.</li>
        <li>Conter pelo menos um numero.</li>
        <li>Conter pelo menos um caractere especial (ex: @, #, $, %, &).</li>
        <li>Proibido utilizar dados obvios como nome do usuario, data de nascimento ou a palavra "cloudscale".</li>
    </ul>

    <h2>2. Autenticacao de Multiplo Fator (2FA)</h2>
    <p>A <strong>autenticacao 2FA e obrigatoria</strong> para todos os acessos de administradores e desenvolvedores aos sistemas internos da CloudScale, incluindo AWS, GCP e painel do cliente B2B. O uso de aplicativos autenticadores (Google Authenticator ou Microsoft Authenticator) e exigido, nao sendo permitido o envio de codigos apenas por SMS por questoes de seguranca.</p>

    <h2>3. Certificacoes e Auditorias</h2>
    <p>A CloudScale Solutions passa por auditorias anuais para manutencao da conformidade com a norma <strong>ISO 27001</strong> (Seguranca da Informacao). Todos os processos internos de desenvolvimento de software seguem as diretrizes do OWASP Top 10.</p>
</body>
</html>
"""
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(content)

def generate_all():
    print(f"Gerando todos os arquivos ficticios em: {DATA_DIR.resolve()}")
    generate_pdf()
    generate_docx()
    generate_xlsx()
    generate_pptx()
    generate_md()
    generate_csv()
    generate_json()
    generate_html()
    print("[OK] Todos os 8 arquivos ficticios da CloudScale foram gerados com sucesso!")

if __name__ == "__main__":
    generate_all()
