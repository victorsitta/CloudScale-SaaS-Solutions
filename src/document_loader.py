import os
import json
from pathlib import Path
from typing import List
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader

class UniversalDocumentLoader:
    """
    Carregador universal capaz de processar os 8 formatos de arquivo exigidos pelo edital
    (PDF, DOCX, XLSX, PPTX, MD, CSV, JSON, HTML) e retornar objetos Document do LangChain.
    """
    
    @staticmethod
    def load_pdf(file_path: Path) -> List[Document]:
        """Carrega arquivos PDF usando o PyPDFLoader do LangChain."""
        loader = PyPDFLoader(str(file_path))
        return loader.load()

    @staticmethod
    def load_docx(file_path: Path) -> List[Document]:
        """Carrega arquivos Word (.docx) usando python-docx."""
        doc = DocxDocument(file_path)
        paragraphs_text = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs_text.append(text)
        
        # Também extrai texto de tabelas dentro do documento Word, se houver
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs_text.append(row_text)
                    
        full_text = "\n".join(paragraphs_text)
        return [Document(
            page_content=full_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "file_type": file_path.suffix.lower()
            }
        )]

    @staticmethod
    def load_xlsx(file_path: Path) -> List[Document]:
        """Carrega planilhas Excel (.xlsx) via pandas e openpyxl."""
        # Carrega todas as abas
        excel_file = pd.ExcelFile(file_path, engine="openpyxl")
        documents = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            sheet_lines = []
            sheet_lines.append(f"Aba/Planilha: {sheet_name}")
            
            for index, row in df.iterrows():
                row_items = []
                for col in df.columns:
                    val = row[col]
                    if pd.notna(val):
                        row_items.append(f"{col}: {val}")
                if row_items:
                    sheet_lines.append(f"Linha {index + 1}: " + " | ".join(row_items))
                    
            full_text = "\n".join(sheet_lines)
            documents.append(Document(
                page_content=full_text,
                metadata={
                    "source": str(file_path),
                    "file_name": file_path.name,
                    "file_type": file_path.suffix.lower(),
                    "sheet_name": sheet_name
                }
            ))
        return documents

    @staticmethod
    def load_csv(file_path: Path) -> List[Document]:
        """Carrega arquivos CSV via pandas."""
        df = pd.read_csv(file_path, encoding="utf-8")
        rows_text = []
        
        for index, row in df.iterrows():
            row_items = []
            for col in df.columns:
                val = row[col]
                if pd.notna(val):
                    row_items.append(f"{col}: {val}")
            if row_items:
                rows_text.append(f"Registro {index + 1}: " + " | ".join(row_items))
                
        full_text = "\n".join(rows_text)
        return [Document(
            page_content=full_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "file_type": file_path.suffix.lower()
            }
        )]

    @staticmethod
    def load_pptx(file_path: Path) -> List[Document]:
        """Carrega apresentações PowerPoint (.pptx) extraindo texto slide por slide."""
        prs = Presentation(file_path)
        documents = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            slide_text_parts.append(f"Slide {slide_idx + 1}")
            
            for shape in slide.shapes:
                # Extrai texto de caixas de texto e formas com texto
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)
                
                # Extrai texto de tabelas nos slides
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_text_parts.append(row_text)
                            
            full_text = "\n".join(slide_text_parts)
            documents.append(Document(
                page_content=full_text,
                metadata={
                    "source": str(file_path),
                    "file_name": file_path.name,
                    "file_type": file_path.suffix.lower(),
                    "slide_number": slide_idx + 1
                }
            ))
        return documents

    @staticmethod
    def load_md(file_path: Path) -> List[Document]:
        """Carrega arquivos Markdown (.md) lendo seu texto bruto."""
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return [Document(
            page_content=content,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "file_type": file_path.suffix.lower()
            }
        )]

    @staticmethod
    def load_html(file_path: Path) -> List[Document]:
        """Carrega arquivos HTML (.html) e extrai o texto limpo usando BeautifulSoup."""
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            
        # Remove tags indesejadas
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
            
        text = soup.get_text(separator="\n")
        # Limpeza de linhas em branco extras
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase for phrase in lines if phrase)
        clean_text = "\n".join(chunks)
        
        title = soup.title.string if soup.title else file_path.name
        
        return [Document(
            page_content=clean_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "file_type": file_path.suffix.lower(),
                "title": title
            }
        )]

    @classmethod
    def load_json(cls, file_path: Path) -> List[Document]:
        """Carrega arquivos JSON (.json) convertendo estruturas aninhadas em blocos de texto."""
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
            
        flat_lines = []
        
        def recurse_json(obj, prefix=""):
            if isinstance(obj, dict):
                for k, v in obj.items():
                    new_prefix = f"{prefix} -> {k}" if prefix else k
                    recurse_json(v, new_prefix)
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    new_prefix = f"{prefix} [{i}]"
                    recurse_json(item, new_prefix)
            else:
                flat_lines.append(f"{prefix}: {obj}")
                
        recurse_json(data)
        full_text = "\n".join(flat_lines)
        
        return [Document(
            page_content=full_text,
            metadata={
                "source": str(file_path),
                "file_name": file_path.name,
                "file_type": file_path.suffix.lower()
            }
        )]

    @classmethod
    def load_file(cls, file_path: Path) -> List[Document]:
        """Mapeia a extensão do arquivo para o parser correto e retorna os Documents."""
        ext = file_path.suffix.lower()
        if not file_path.exists():
            raise FileNotFoundError(f"Arquivo não encontrado: {file_path}")
            
        if ext == ".pdf":
            docs = cls.load_pdf(file_path)
        elif ext == ".docx":
            docs = cls.load_docx(file_path)
        elif ext == ".xlsx":
            docs = cls.load_xlsx(file_path)
        elif ext == ".csv":
            docs = cls.load_csv(file_path)
        elif ext == ".pptx":
            docs = cls.load_pptx(file_path)
        elif ext == ".md":
            docs = cls.load_md(file_path)
        elif ext in [".html", ".htm"]:
            docs = cls.load_html(file_path)
        elif ext == ".json":
            docs = cls.load_json(file_path)
        else:
            # Fallback para leitura como texto simples
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                docs = [Document(
                    page_content=content,
                    metadata={
                        "source": str(file_path),
                        "file_name": file_path.name,
                        "file_type": ext
                    }
                )]
            except Exception as e:
                raise ValueError(f"Extensão de arquivo não suportada: {ext}. Erro fallback: {str(e)}")
                
        # Garante que os metadados básicos estejam presentes em todos os documentos
        for d in docs:
            d.metadata.setdefault("source", str(file_path))
            d.metadata.setdefault("file_name", file_path.name)
            d.metadata.setdefault("file_type", ext)
            
        return docs

    @classmethod
    def load_directory(cls, dir_path: Path) -> List[Document]:
        """Varre um diretório carregando todos os arquivos suportados."""
        all_docs = []
        if not dir_path.exists():
            return []
            
        for root, _, files in os.walk(dir_path):
            for file in files:
                file_path = Path(root) / file
                # Ignora arquivos ocultos ou temporários
                if file.startswith(".") or file.startswith("~$"):
                    continue
                try:
                    docs = cls.load_file(file_path)
                    all_docs.extend(docs)
                except Exception as e:
                    print(f"Erro ao carregar o arquivo {file_path}: {e}")
                    
        return all_docs
